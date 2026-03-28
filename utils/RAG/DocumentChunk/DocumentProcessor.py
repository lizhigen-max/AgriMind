#!/usr/bin/env python
# -*- coding: UTF-8 -*-
'''
@Project ：AgriMind
@File    ：DocumentProcessor.py
@IDE     ：PyCharm
@Author  ：lizhigen
@Email   ：lizhigen1996@aliyun.com
@Date    ：2026/3/28 21:08
'''

import json
import os
from typing import List, Union, Optional, Dict, Tuple
from langchain_core.documents import Document
from langchain_text_splitters import (
    RecursiveCharacterTextSplitter,
    CharacterTextSplitter,
    MarkdownTextSplitter,
    TokenTextSplitter
)
from langchain_community.document_loaders import (
    TextLoader,
    PyPDFLoader,
    Docx2txtLoader,
    UnstructuredMarkdownLoader,
    CSVLoader
)
import pandas as pd
import pptx
import re
import random
import warnings
import logging
from .WordImage import word_extract_and_replace_images

warnings.filterwarnings('ignore')


class DocumentProcessor:
    """
    多功能文档处理器，支持多种格式文档的加载与分割

    支持格式：
    - txt, md, csv: 文本类
    - pdf: PDF文档
    - doc, docx: Word文档
    - xls, xlsx: Excel表格
    - ppt, pptx: PowerPoint演示文稿
    """

    def __init__(self,
                 chunk_size: int = 500,
                 chunk_overlap: int = 50,
                 splitter_type: str = "recursive",
                 encoding: str = "utf-8",
                 clean_config: Optional[Dict] = None):
        """
        初始化文档处理器

        Args:
            chunk_size: 每个文本块的最大字符数
            chunk_overlap: 文本块之间的重叠字符数
            splitter_type: 分割器类型，可选 "recursive", "character", "markdown", "token"
            encoding: 文件编码
            clean_config: 清洗配置字典
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.splitter_type = splitter_type
        self.encoding = encoding

        # 初始化文本分割器
        self.splitter = self._init_splitter()

        # 设置清洗配置
        self.clean_config = clean_config or {
            'remove_extra_whitespace': True,  #去除多余空白字符
            'remove_special_chars': False,  # 去除特殊字符
            'strip_whitespace': True,  # 去除首尾空白
            'lowercase': False,  # 转换为小写
            'remove_urls': True,  # 去除url
            'remove_emails': True,  # 去除电子邮件
            'remove_phone_numbers': True,  # 去除电话号码
            'remove_numbers': False,  # 去除纯数字
            'remove_punctuation': False,  # 去除标点符号
            'custom_patterns': None,  # 应用自定义正则表达式模式
            'min_word_length': 1,  # 过滤过短或过长的单词
            'max_word_length': 50,
            'stop_words': None  # 去除停用词
        }

        # 预编译常用正则表达式
        self.patterns = {
            'extra_spaces': re.compile(r'\s+'),
            'special_chars': re.compile(r'[^\w\s\u4e00-\u9fff]'),  # 保留中文字符
            'numbers': re.compile(r'\b\d+\b'),
            'page_header': re.compile(r'^第[一二三四五六七八九十\d]+页\s*$', re.MULTILINE),
            'page_footer': re.compile(r'^\s*-\s*\d+\s*-\s*$', re.MULTILINE),
            'repeated_chars': re.compile(r'(.)\1{3,}'),  # 重复4次以上的字符
        }

        # 根据文件类型调用对应的加载器
        self.loaders = {
            '.txt': self._load_txt,
            '.pdf': self._load_pdf,
            '.docx': self._load_word,
            '.doc': self._load_word,  # 旧版Word文档
            '.xlsx': self._load_excel,
            '.xls': self._load_excel,
            '.csv': self._load_csv,
            '.pptx': self._load_ppt,
            '.ppt': self._load_ppt,
            '.md': self._load_markdown,
            '.markdown': self._load_markdown
        }

    def _init_splitter(self):
        """初始化文本分割器"""
        if self.splitter_type == "recursive":
            return RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                length_function=len,
                separators=["\n\n", "\n", "。", "；", "，", " ", ""]
            )
        elif self.splitter_type == "character":
            return CharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                separator="\n"
            )
        elif self.splitter_type == "markdown":
            return MarkdownTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap
            )
        elif self.splitter_type == "token":
            return TokenTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap
            )
        else:
            raise ValueError(f"不支持的splitter_type: {self.splitter_type}")

    def clean_text(self, text: str, doc_type: str = None) -> str:
        """
        清洗文本数据

        Args:
            text: 原始文本
            doc_type: 文档类型，用于特定类型的清洗

        Returns:
            str: 清洗后的文本
        """
        if not text or not text.strip():
            return ""

        cleaned_text = text

        # 1. 去除URL
        if self.clean_config.get('remove_urls', True):
            patterns = [
                # http:// 或 https:// 开头的URL
                r'https?://(?:[a-zA-Z0-9-]+\.)+[a-zA-Z]{2,}(?:/[^\s\u4e00-\u9fff<>"\']*)?',
                # www. 开头的URL
                r'www\.(?:[a-zA-Z0-9-]+\.)+[a-zA-Z]{2,}(?:/[^\s\u4e00-\u9fff<>"\']*)?',
            ]

            for pattern in patterns:
                cleaned_text = re.sub(pattern, '', cleaned_text)

            # 清理多余空格，但保留原有结构
            cleaned_text = re.sub(r'\s{2,}', ' ', cleaned_text)
            cleaned_text = re.sub(r'^\s+|\s+$', '', cleaned_text)

        # 2. 去除电子邮件
        if self.clean_config.get('remove_emails', True):
            pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b'
            cleaned_text =  re.sub(pattern, '', cleaned_text, flags=re.ASCII)

        # 3. 去除电话号码
        if self.clean_config.get('remove_phone_numbers', True):
            mobile_pattern = r'(?<!\d)1[3-9]\d[- .]?\d{4}[- .]?\d{4}(?!\d)'
            landline_pattern = r'(?<!\d)(?:0\d{2,3}[- .])?\d{7,8}(?!\d)'
            pattern = f'{mobile_pattern}|{landline_pattern}'
            cleaned_text = re.sub(pattern, '', cleaned_text)

        # 4. 去除多余空白字符
        if self.clean_config.get('remove_extra_whitespace', True):
            cleaned_text = self.patterns['extra_spaces'].sub(' ', cleaned_text)

        # 5. 去除页眉页脚（针对PDF/DOC）
        if doc_type in ['pdf', 'doc', 'docx']:
            cleaned_text = self.patterns['page_header'].sub('', cleaned_text)
            cleaned_text = self.patterns['page_footer'].sub('', cleaned_text)

        # 6. 去除重复字符（如"------"或"******"）
        cleaned_text = self.patterns['repeated_chars'].sub(r'\1\1\1', cleaned_text)

        # 7. 去除特殊字符
        if self.clean_config.get('remove_special_chars', False):
            cleaned_text = self.patterns['special_chars'].sub(' ', cleaned_text)

        # 8. 去除纯数字
        if self.clean_config.get('remove_numbers', False):
            cleaned_text = self.patterns['numbers'].sub('', cleaned_text)

        # 9. 转换为小写
        if self.clean_config.get('lowercase', False):
            cleaned_text = cleaned_text.lower()

        # 10. 去除标点符号
        if self.clean_config.get('remove_punctuation', False):
            import string
            translator = str.maketrans('', '', string.punctuation)
            cleaned_text = cleaned_text.translate(translator)

        # 11. 应用自定义正则表达式模式
        custom_patterns = self.clean_config.get('custom_patterns')
        if custom_patterns:
            for pattern, replacement in custom_patterns.items():
                cleaned_text = re.compile(pattern).sub(replacement, cleaned_text)

        # 12. 过滤过短或过长的单词
        # min_len = self.clean_config.get('min_word_length', 1)
        # max_len = self.clean_config.get('max_word_length', 50)
        # words = cleaned_text.split()
        # filtered_words = []
        # for word in words:
        #     if min_len <= len(word) <= max_len:
        #         filtered_words.append(word)
        # cleaned_text = ' '.join(filtered_words)

        # 13. 去除停用词
        stop_words = self.clean_config.get('stop_words')
        if stop_words:
            words = cleaned_text.split()
            filtered_words = [word for word in words if word.lower() not in stop_words]
            cleaned_text = ' '.join(filtered_words)

        # 14. 去除首尾空白
        if self.clean_config.get('strip_whitespace', True):
            cleaned_text = cleaned_text.strip()

        return cleaned_text

    def clean_document(self, document: Document, doc_type: str = None) -> Document:
        """
        清洗单个文档

        Args:
            document: 文档对象
            doc_type: 文档类型

        Returns:
            Document: 清洗后的文档
        """
        cleaned_content = self.clean_text(document.page_content, doc_type)

        # 如果清洗后内容为空，返回None
        if not cleaned_content.strip():
            return None

        # 创建新的文档对象，保留原始元数据
        cleaned_doc = Document(
            page_content=cleaned_content,
            metadata=document.metadata.copy()
        )

        # 在元数据中添加清洗标记
        cleaned_doc.metadata['cleaned'] = True
        cleaned_doc.metadata['original_length'] = len(document.page_content)
        cleaned_doc.metadata['cleaned_length'] = len(cleaned_content)

        return cleaned_doc

    def clean_documents(self, documents: List[Document], doc_type: str = None) -> List[Document]:
        """
        批量清洗文档

        Args:
            documents: 文档列表
            doc_type: 文档类型

        Returns:
            List[Document]: 清洗后的文档列表
        """
        cleaned_docs = []
        for doc in documents:
            cleaned_doc = self.clean_document(doc, doc_type)
            if cleaned_doc:
                cleaned_docs.append(cleaned_doc)

        return cleaned_docs

    def _load_txt(self, file_path: str, is_ocr: bool) -> List[Document]:
        """加载文本文件"""
        try:
            loader = TextLoader(file_path, encoding=self.encoding)
            documents = loader.load()
            return documents
        except Exception as e:
            # 备用方案：直接读取文件
            try:
                with open(file_path, 'r', encoding=self.encoding) as f:
                    content = f.read()
                doc = Document(page_content=content, metadata={"source": file_path, "type": "txt"})
                return [doc] if doc else []
            except Exception as e2:
                raise Exception(f"文本文件加载失败: {str(e2)}")

    def _load_pdf(self, file_path: str, is_ocr: bool) -> List[Document]:
        """加载PDF文件"""
        try:
            loader = PyPDFLoader(file_path)
            return loader.load()
        except Exception as e:
            # 尝试使用备用PDF加载器
            try:
                from langchain_community.document_loaders import UnstructuredPDFLoader
                loader = UnstructuredPDFLoader(file_path, mode="elements")
                return loader.load()
            except:
                raise Exception(f"PDF加载失败: {str(e)}")

    def _load_word(self, file_path: str, is_ocr: bool) -> List[Document]:
        """加载Word文档"""
        try:
            # 处理图片
            images, new_file_path = word_extract_and_replace_images(file_path, is_ocr=is_ocr)
            loader = Docx2txtLoader(new_file_path)
            return loader.load()
        except Exception as e:
            # 尝试使用备用Word加载器
            try:
                from langchain_community.document_loaders import UnstructuredWordDocumentLoader
                loader = UnstructuredWordDocumentLoader(file_path)
                return loader.load()
            except:
                raise Exception(f"Word文档加载失败: {str(e)}")

    def _load_excel(self, file_path: str, is_ocr: bool) -> List[Document]:
        """加载Excel文件"""
        try:
            # 读取Excel文件
            if file_path.endswith('.xlsx'):
                df = pd.read_excel(file_path, engine='openpyxl')
            elif file_path.endswith('.xls'):
                df = pd.read_excel(file_path, engine='xlrd')
            else:
                raise ValueError("不支持的Excel格式")

            # 将DataFrame转换为文本
            documents = []
            for idx, row in df.iterrows():
                # 将每行数据转换为文本
                content_parts = []
                for col, value in row.items():
                    if pd.notna(value):
                        content_parts.append(f"{col}: {value}")

                if content_parts:
                    text_content = f"行 {idx + 1}: " + "; ".join(content_parts)
                    metadata = {
                        "source": file_path,
                        "row": idx + 1,
                        "type": "excel"
                    }
                    documents.append(Document(page_content=text_content, metadata=metadata))

            return documents
        except Exception as e:
            raise Exception(f"Excel文件加载失败: {str(e)}")

    def _load_csv(self, file_path: str, is_ocr: bool) -> List[Document]:
        """加载CSV文件"""
        try:
            loader = CSVLoader(file_path, encoding=self.encoding)
            return loader.load()
        except Exception as e:
            # 如果标准CSV加载器失败，尝试使用pandas
            try:
                df = pd.read_csv(file_path, encoding=self.encoding)
                documents = []
                for idx, row in df.iterrows():
                    content_parts = [f"{col}: {val}" for col, val in row.items() if pd.notna(val)]
                    if content_parts:
                        text_content = f"行 {idx + 1}: " + "; ".join(content_parts)
                        metadata = {"source": file_path, "row": idx + 1, "type": "csv"}
                        documents.append(Document(page_content=text_content, metadata=metadata))
                return documents
            except:
                raise Exception(f"CSV文件加载失败: {str(e)}")

    def _load_ppt(self, file_path: str, is_ocr: bool) -> List[Document]:
        """加载PowerPoint文件"""
        try:
            presentation = pptx.Presentation(file_path)
            documents = []

            for slide_idx, slide in enumerate(presentation.slides):
                slide_text_parts = []

                # 提取幻灯片标题
                if slide.shapes.title:
                    slide_text_parts.append(f"标题: {slide.shapes.title.text}")

                # 提取所有形状的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if not (shape == slide.shapes.title):
                            slide_text_parts.append(shape.text)

                if slide_text_parts:
                    text_content = "\n".join(slide_text_parts)
                    metadata = {
                        "source": file_path,
                        "slide": slide_idx + 1,
                        "total_slides": len(presentation.slides),
                        "type": "ppt"
                    }
                    documents.append(Document(page_content=text_content, metadata=metadata))

            return documents
        except Exception as e:
            raise Exception(f"PPT文件加载失败: {str(e)}。请确保已安装python-pptx: pip install python-pptx")

    def _load_markdown(self, file_path: str, is_ocr: bool) -> List[Document]:
        """加载Markdown文件"""
        try:
            loader = UnstructuredMarkdownLoader(file_path)
            return loader.load()
        except Exception as e:
            # 备用方案
            try:
                with open(file_path, 'r', encoding=self.encoding) as f:
                    content = f.read()
                return [Document(page_content=content, metadata={"source": file_path})]
            except:
                raise Exception(f"Markdown文件加载失败: {str(e)}")

    def load_document(self, file_path: str, is_ocr: bool) -> Tuple[List[Document], str]:
        """
        根据文件扩展名加载文档

        Args:
            file_path: 文档路径

        Returns:
            List[Document]: 加载的文档列表
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        # 获取文件扩展名
        file_ext = os.path.splitext(file_path)[1].lower()

        if file_ext in self.loaders:
            logging.info(f"正在加载 {file_ext} 文件: {os.path.basename(file_path)}")
            # 结合数据清洗
            return self.clean_documents(self.loaders[file_ext](file_path, is_ocr), file_ext), file_path
        else:
            supported_formats = ', '.join(self.loaders.keys())
            raise ValueError(
                f"不支持的文件格式: {file_ext}\n"
                f"当前支持的格式: {supported_formats}\n"
                f"请转换为支持的格式或使用其他工具处理。"
            )

    def split_document(self, file_path: str, is_ocr: bool, **splitter_kwargs) -> List[Document]:
        """
        加载并分割文档

        Args:
            file_path: 文档路径
            **splitter_kwargs: 分割器参数，可覆盖默认参数

        Returns:
            List[Document]: 分割后的文档块列表
        """
        # 加载文档
        documents, orignal_file_path = self.load_document(file_path, is_ocr)

        if not documents:
            logging.info(f"警告: 文件 {file_path} 没有加载到任何内容")
            return []

        # 更新分割器参数（如果提供）
        if splitter_kwargs:
            for key, value in splitter_kwargs.items():
                if hasattr(self.splitter, key):
                    setattr(self.splitter, key, value)
                else:
                    logging.info(f"警告: 分割器没有参数 {key}，已忽略")

        # 分割文档
        logging.info(f"正在分割文档，原始文档数: {len(documents)}")
        chunks = self.splitter.split_documents(documents)
        for chunk in chunks:
            chunk.metadata.update({"orignal_file_path": orignal_file_path})
        logging.info(f"分割完成，生成 {len(chunks)} 个文本块")

        return chunks

    def batch_process(self,
                      floder: str,
                      show_progress: bool = True, is_ocr=False) -> List[Document]:
        """
        批量处理多个文档
        自动为分块添加标签

        Args:
            floder: 文档保存的文件夹路径列表
            show_progress: 是否显示处理进度
            is_ocr: word是否启用ocr文字识别， 暂时只只支持word

        Returns:
            List[Document]: 所有分割后的文档块
        """
        if not os.path.exists(floder):
            logging.info(f'文件夹{floder}不存在')
            return []
        tag_path = floder + os.sep + 'DocumentTag.json'
        documentTag = []
        if os.path.exists(tag_path):
            with open(tag_path, 'r', encoding='utf-8') as fp:
                documentTag = json.loads(fp.read())

        file_paths = []
        for file in os.listdir(floder):
            if file == 'DocumentTag.json':
                # 跳过标签文件
                continue
            # 获取文件扩展名
            file_path = floder + os.sep + file
            if os.path.isdir(file_path):
                continue
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext not in self.loaders:
                logging.info(f'当前暂不支持处理{file_path}类型文件')
                continue
            isFound = False
            for tag in documentTag:
                if tag['file_name'] == file:
                    isFound = True
                    file_paths.append([file_path, tag])
                    break
            if not isFound:
                file_paths.append([file_path, {}])

        all_chunks = []
        total_files = len(file_paths)

        for idx, (file_path, tag) in enumerate(file_paths, 1):
            if show_progress:
                logging.info(f"\n处理文件 {idx}/{total_files}: {os.path.basename(file_path)}")

            try:
                chunks = self.split_document(file_path, is_ocr)
                for chunk in chunks:
                    chunk.metadata.update(tag)  # 为分块添加标签
                all_chunks.extend(chunks)
                chunk = random.choice(chunks)
                logging.info('随机块内容: ' + str(chunk.page_content))
                logging.info('随机块信息: ' + str(chunk.metadata))

                if show_progress:
                    logging.info(f"✓ 成功处理，生成 {len(chunks)} 个文本块")

            except Exception as e:
                logging.info(f"✗ 处理失败: {str(e)}")
                continue

        if show_progress:
            logging.info(f"\n批量处理完成，总共生成 {len(all_chunks)} 个文本块")

        return all_chunks

    def get_statistics(self, chunks: List[Document]) -> dict:
        """
        获取文档块统计信息

        Args:
            chunks: 文档块列表

        Returns:
            dict: 统计信息
        """
        if not chunks:
            return {}

        # 计算文本块长度
        chunk_lengths = [len(chunk.page_content) for chunk in chunks]

        # 提取元数据中的文件类型
        file_types = {}
        for chunk in chunks:
            doc_type = chunk.metadata.get("type", "unknown")
            file_types[doc_type] = file_types.get(doc_type, 0) + 1

        return {
            "total_chunks": len(chunks),
            "avg_chunk_length": sum(chunk_lengths) / len(chunk_lengths),
            "min_chunk_length": min(chunk_lengths),
            "max_chunk_length": max(chunk_lengths),
            "file_types": file_types
        }

    def export_chunks(self,
                      chunks: List[Document],
                      output_file: str = None,
                      format: str = "txt") -> Optional[str]:
        """
        导出分割后的文本块

        Args:
            chunks: 文档块列表
            output_file: 输出文件路径，如果为None则返回字符串
            format: 输出格式，支持 "txt" 或 "json"

        Returns:
            Optional[str]: 如果output_file为None，则返回字符串
        """
        if format == "txt":
            content_lines = []
            for i, chunk in enumerate(chunks, 1):
                content_lines.append(f"=== 文本块 {i} ===")
                content_lines.append(f"长度: {len(chunk.page_content)} 字符")
                content_lines.append(f"元数据: {chunk.metadata}")
                content_lines.append(chunk.page_content)
                content_lines.append("=" * 50 + "\n")

            content = "\n".join(content_lines)

        elif format == "json":
            import json
            chunks_data = []
            for chunk in chunks:
                chunks_data.append({
                    "content": chunk.page_content,
                    "metadata": chunk.metadata
                })
            content = json.dumps(chunks_data, ensure_ascii=False, indent=2)
        else:
            raise ValueError(f"不支持的格式: {format}，支持 'txt' 或 'json'")

        if output_file:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(content)
            logging.info(f"导出完成: {output_file}")
            return None
        else:
            return content


# 使用示例
if __name__ == "__main__":
    # 初始化处理器
    processor = DocumentProcessor(
        chunk_size=500,
        chunk_overlap=50,
        splitter_type="recursive"
    )

    # # 示例1: 处理单个文件
    # try:
    #     # 请替换为实际文件路径
    #     chunks = processor.split_document('../Document/基于气象及虫情监测数据的茶小绿叶蝉虫害动态预测方法.pdf')
    #
    #     # 获取统计信息
    #     stats = processor.get_statistics(chunks)
    #     logging.info("统计信息:", stats)
    #     chunk = random.choice(chunks)
    #     logging.info('随机块内容: ', chunk.page_content)
    #     logging.info('随机块信息: ', chunk.metadata)
    #
    #     # 导出结果
    #     # processor.export_chunks(chunks, "output.txt", format="txt")
    #     pass
    #
    # except Exception as e:
    #     logging.info(f"处理失败: {e}")
    processor.batch_process('../Document', show_progress=True)
